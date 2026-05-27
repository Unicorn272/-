import io
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# 색상
NAVY = RGBColor(0x1B, 0x3A, 0x6B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x88, 0x88, 0x88)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ── 헬퍼 ──────────────────────────────────────────────────────────────

def _add_header(slide, title: str):
    """네이비 헤더 바 + 흰 제목"""
    bar = slide.shapes.add_shape(1, 0, 0, SLIDE_W, Inches(1.1))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()

    tf = bar.text_frame
    tf.text = title
    tf.paragraphs[0].runs[0].font.color.rgb = WHITE
    tf.paragraphs[0].runs[0].font.size = Pt(24)
    tf.paragraphs[0].runs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.LEFT
    tf.margin_left = Inches(0.4)
    tf.margin_top = Inches(0.25)


def _add_source_bar(slide, sources: list[dict]):
    """
    슬라이드 하단 출처 바.
    sources: [{"text": "삼성전기 증권신고서 2025.03", "type": "filing"}, ...]
    """
    bar = slide.shapes.add_shape(1, 0, Inches(7.0), SLIDE_W, Inches(0.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = WHITE
    bar.line.fill.background()

    tf = bar.text_frame
    tf.margin_left = Inches(0.3)
    tf.margin_top = Inches(0.05)

    para = tf.paragraphs[0]
    for i, src in enumerate(sources):
        if i > 0:
            run = para.add_run()
            run.text = "  |  "
            run.font.size = Pt(8)
            run.font.color.rgb = GRAY

        run = para.add_run()
        label = "[신고서]" if src["type"] == "filing" else "[뉴스]"
        run.text = f"{label} {src['text']}"
        run.font.size = Pt(8)
        run.font.color.rgb = NAVY if src["type"] == "filing" else GRAY


def _add_text_box(slide, text, left, top, width, height, size=12, bold=False, color=None):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf = txb.text_frame
    tf.word_wrap = True
    tf.text = text
    run = tf.paragraphs[0].runs[0]
    run.font.size = Pt(size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color


def _radar_chart(labels: list[str], scores: list[int]) -> bytes:
    """5 Forces 레이더 차트 → PNG bytes"""
    N = len(labels)
    angles = np.linspace(0, 2 * np.pi, N, endpoint=False).tolist()
    angles += angles[:1]
    values = scores + scores[:1]

    fig, ax = plt.subplots(figsize=(4, 4), subplot_kw=dict(polar=True))
    ax.plot(angles, values, color="#1B3A6B", linewidth=2)
    ax.fill(angles, values, color="#1B3A6B", alpha=0.2)
    ax.set_xticks(angles[:-1])
    ax.set_xticklabels(labels, fontsize=9)
    ax.set_ylim(0, 5)
    ax.set_yticks([1, 2, 3, 4, 5])
    ax.set_yticklabels(["1", "2", "3", "4", "5"], fontsize=7)
    fig.tight_layout()

    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=150)
    plt.close(fig)
    buf.seek(0)
    return buf.read()


# ── 슬라이드 생성 ──────────────────────────────────────────────────────

def _slide_cover(prs, industry: str, analysis_date: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    bg = slide.shapes.add_shape(1, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()

    _add_text_box(slide, industry, Inches(1), Inches(2.5), Inches(10), Inches(1.5),
                  size=36, bold=True, color=WHITE)
    _add_text_box(slide, "산업분석 보고서", Inches(1), Inches(3.8), Inches(8), Inches(0.8),
                  size=18, color=WHITE)
    _add_text_box(slide, analysis_date, Inches(1), Inches(6.5), Inches(6), Inches(0.5),
                  size=11, color=GRAY)


def _slide_executive_summary(prs, insights: list[str], sources: list[dict]):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_header(slide, "Executive Summary")

    for i, insight in enumerate(insights[:5]):
        top = Inches(1.4 + i * 0.95)
        bullet = slide.shapes.add_shape(1, Inches(0.4), top + Inches(0.1),
                                         Inches(0.08), Inches(0.08))
        bullet.fill.solid()
        bullet.fill.fore_color.rgb = NAVY
        bullet.line.fill.background()

        _add_text_box(slide, insight, Inches(0.65), top, Inches(11.5), Inches(0.85), size=13)

    _add_source_bar(slide, sources)


def _slide_forces_summary(prs, five_forces: dict, sources: list[dict]):
    labels = ["경쟁강도", "공급자\n협상력", "구매자\n협상력", "대체재\n위협", "신규진입\n위협"]
    scores = [
        five_forces["competitive_rivalry"]["score"],
        five_forces["supplier_power"]["score"],
        five_forces["buyer_power"]["score"],
        five_forces["threat_of_substitutes"]["score"],
        five_forces["threat_of_new_entrants"]["score"],
    ]

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_header(slide, "Porter 5 Forces — 종합")

    chart_bytes = _radar_chart(labels, scores)
    chart_buf = io.BytesIO(chart_bytes)
    slide.shapes.add_picture(chart_buf, Inches(4.2), Inches(1.3), Inches(4.5), Inches(4.5))

    # 점수 요약 텍스트
    summary_lines = "\n".join(
        f"{l.replace(chr(10), ' ')}  {'●' * s}{'○' * (5 - s)}  ({s}/5)"
        for l, s in zip(labels, scores)
    )
    _add_text_box(slide, summary_lines, Inches(0.4), Inches(1.4), Inches(3.6), Inches(5),
                  size=11)

    _add_source_bar(slide, sources)


def _slide_force_detail(prs, force_name: str, data: dict, sources: list[dict]):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_header(slide, f"Porter 5 Forces — {force_name}")

    score_bar = "●" * data["score"] + "○" * (5 - data["score"])
    _add_text_box(slide, f"위험도  {score_bar}  ({data['score']}/5)",
                  Inches(0.4), Inches(1.3), Inches(6), Inches(0.5), size=13, bold=True)

    _add_text_box(slide, data["summary"],
                  Inches(0.4), Inches(1.9), Inches(12), Inches(1.0), size=12)

    _add_text_box(slide, "근거", Inches(0.4), Inches(3.1), Inches(2), Inches(0.4),
                  size=11, bold=True, color=NAVY)

    evidence_text = "\n".join(f"• {e}" for e in data.get("evidence", []))
    _add_text_box(slide, evidence_text, Inches(0.4), Inches(3.55), Inches(12), Inches(3.0),
                  size=10)

    _add_source_bar(slide, sources)


def _slide_pestel_summary(prs, pestel: dict, sources: list[dict]):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_header(slide, "PESTEL 분석 — 종합")

    items = [
        ("Political", pestel["political"]["summary"]),
        ("Economic", pestel["economic"]["summary"]),
        ("Social", pestel["social"]["summary"]),
        ("Technological", pestel["technological"]["summary"]),
        ("Environmental", pestel["environmental"]["summary"]),
        ("Legal", pestel["legal"]["summary"]),
    ]

    for idx, (label, summary) in enumerate(items):
        col = idx % 3
        row = idx // 3
        left = Inches(0.3 + col * 4.3)
        top = Inches(1.3 + row * 2.7)

        box = slide.shapes.add_shape(1, left, top, Inches(4.0), Inches(2.4))
        box.fill.solid()
        box.fill.fore_color.rgb = WHITE
        box.line.color.rgb = NAVY

        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.1)
        tf.margin_top = Inches(0.08)

        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = label
        r.font.bold = True
        r.font.size = Pt(12)
        r.font.color.rgb = NAVY

        p2 = tf.add_paragraph()
        r2 = p2.add_run()
        r2.text = summary
        r2.font.size = Pt(10)

    _add_source_bar(slide, sources)


def _slide_pestel_detail(prs, pestel: dict, sources: list[dict]):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_header(slide, "PESTEL 분석 — 상세")

    keys = ["political", "economic", "social", "technological", "environmental", "legal"]
    labels = ["P", "E", "S", "T", "E", "L"]

    for idx, (key, label) in enumerate(zip(keys, labels)):
        col = idx % 3
        row = idx // 3
        left = Inches(0.3 + col * 4.3)
        top = Inches(1.3 + row * 2.7)

        box = slide.shapes.add_shape(1, left, top, Inches(4.0), Inches(2.4))
        box.fill.solid()
        box.fill.fore_color.rgb = WHITE
        box.line.color.rgb = NAVY

        item = pestel[key]
        src_tag = f"[{item.get('source_type', '')}]"

        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.1)
        tf.margin_top = Inches(0.08)

        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = f"{label}  {src_tag}"
        r.font.bold = True
        r.font.size = Pt(11)
        r.font.color.rgb = NAVY

        evidence_text = "  ".join(f"• {e}" for e in item.get("evidence", [])[:2])
        p2 = tf.add_paragraph()
        r2 = p2.add_run()
        r2.text = evidence_text
        r2.font.size = Pt(9)

    _add_source_bar(slide, sources)


def _slide_limitations(prs, limitations: list[str]):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_header(slide, "데이터 한계 및 유의사항")

    for i, item in enumerate(limitations):
        top = Inches(1.5 + i * 0.8)
        _add_text_box(slide, f"• {item}", Inches(0.5), top, Inches(12), Inches(0.7), size=12)


# ── 메인 ──────────────────────────────────────────────────────────────

def generate_ppt(analysis: dict, industry: str, analysis_date: str) -> bytes:
    """
    analysis: Agent 3 출력 JSON
    반환: PPT 파일 bytes
    """
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    ff = analysis["five_forces"]
    pestel = analysis["pestel"]

    # 공통 출처 (신고서 기반)
    filing_sources = [{"text": src, "type": "filing"}
                      for src in _collect_sources(ff, pestel, source_type="신고서")]
    news_sources = [{"text": src, "type": "news"}
                    for src in _collect_sources(ff, pestel, source_type="뉴스")]
    all_sources = filing_sources + news_sources

    force_map = {
        "경쟁강도": "competitive_rivalry",
        "공급자 협상력": "supplier_power",
        "구매자 협상력": "buyer_power",
        "대체재 위협": "threat_of_substitutes",
        "신규진입 위협": "threat_of_new_entrants",
    }

    _slide_cover(prs, industry, analysis_date)
    _slide_executive_summary(prs, analysis.get("key_insights", []), all_sources)
    _slide_forces_summary(prs, ff, filing_sources)

    for name, key in force_map.items():
        sources = [{"text": e, "type": "filing"} for e in ff[key].get("evidence", [])]
        _slide_force_detail(prs, name, ff[key], sources)

    _slide_pestel_summary(prs, pestel, all_sources)
    _slide_pestel_detail(prs, pestel, all_sources)
    _slide_limitations(prs, analysis.get("data_limitations", []))

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()


def _collect_sources(ff: dict, pestel: dict, source_type: str) -> list[str]:
    sources = set()
    for force in ff.values():
        for e in force.get("evidence", []):
            if source_type in e:
                sources.add(e)
    for item in pestel.values():
        if item.get("source_type") == source_type:
            for e in item.get("evidence", []):
                sources.add(e)
    return list(sources)
